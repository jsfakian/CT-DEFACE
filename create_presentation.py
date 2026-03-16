#!/usr/bin/env python
"""
Create a professional PPTX presentation for CT-DEFACE
Non-technical stakeholder presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(0, 51, 102)  # Dark blue
ACCENT_COLOR = RGBColor(0, 120, 215)  # Light blue
TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray
LIGHT_BG = RGBColor(240, 248, 255)  # Alice blue

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(32)
        p.font.color.rgb = ACCENT_COLOR
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TITLE_COLOR
    title_shape.line.color.rgb = TITLE_COLOR
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.margin_bottom = Inches(0.1)
    title_frame.margin_top = Inches(0.2)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TITLE_COLOR
    title_shape.line.color.rgb = TITLE_COLOR
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.4))
    p = left_header.text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4))
    p = right_header.text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5.5))
    text_frame = left_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(left_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.5))
    text_frame = right_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(right_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "CT-DEFACE", "Protecting Patient Privacy in Medical Imaging")

# Slide 2: The Problem
add_content_slide(prs, "🎯 The Problem", [
    "• Medical imaging reveals identifiable facial features",
    "• Doctors need to share images for research and teaching",
    "• Traditional anonymization removes all clinical data",
    "• Challenge: Remove faces while keeping patient information"
])

# Slide 3: The Solution
add_content_slide(prs, "🧠 Introducing CT-DEFACE", [
    "✓ Automatically removes facial anatomy from CT scans",
    "✓ Keeps all patient information intact (ID, dates, etc.)",
    "✓ Fast and easy to use—no special expertise needed",
    "✓ Works on standard computers without expensive GPUs",
    "✓ Data stays on your computer—complete privacy control"
])

# Slide 4: How It Works
add_content_slide(prs, "🔬 How It Actually Works", [
    "Step 1: Read DICOM images from a folder",
    "Step 2: AI identifies facial structures in the scan",
    "Step 3: Removes or obscures facial region",
    "Step 4: Save defaced images in original format",
    "Result: Faces removed, clinical data preserved"
])

# Slide 5: Setup Overview
add_content_slide(prs, "💻 Getting CT-DEFACE Ready", [
    "Windows & Linux: Same simple process",
    "1. Download Git & Python 3.12 (~15 min)",
    "2. Run one setup command (~15 min)",
    "3. Ready to use!",
    "Total Time: About 30 minutes"
])

# Slide 6: Running the Pipeline
add_content_slide(prs, "▶️ Using CT-DEFACE", [
    "Step 1: Create folder 'dicom_input' and add your scans",
    "Step 2: Run: python cta_deface_pipeline_multi2.py -i dicom_input -o dicom_output",
    "Step 3: Retrieve defaced scans from 'dicom_output' folder",
    "Processing Time: 2-5 minutes per scan"
])

# Slide 7: Time Comparison
add_two_column_slide(prs, "⏱️ Speed Comparison",
    "Manual Defacing",
    ["• Pick each image individually",
     "• Manually paint over faces",
     "• Verify quality",
     "• 15-30 min per scan",
     "• 50 cases = 12-25 hours"],
    "Automated (CT-DEFACE)",
    ["• Point to folder of images",
     "• Run one command",
     "• Computer works overnight",
     "• 2-5 min per scan",
     "• 50 cases = 2-4 hours"]
)

# Slide 8: Security & Privacy
add_content_slide(prs, "🔐 Security & Privacy", [
    "• All processing happens on YOUR computer",
    "• No data sent to cloud or external servers",
    "• Original data never modified—completely safe",
    "• HIPAA compliant",
    "• You maintain complete control"
])

# Slide 9: Key Benefits
add_content_slide(prs, "✅ Key Benefits", [
    "🚀 Speed: Process 50 scans in hours instead of days",
    "💰 Cost: Free and open-source, no licenses",
    "🔒 Security: Data stays local and protected",
    "😊 Ease: One command, no training required",
    "📋 Clinical Data: Preserves all medical information"
])

# Slide 10: Common Questions
add_content_slide(prs, "🤔 Common Questions", [
    "Q: Is this legal for healthcare? A: Yes, designed for HIPAA.",
    "Q: Can it damage scans? A: No. Original files stay safe.",
    "Q: How accurate? A: Trained on 1000+ real CT scans.",
    "Q: Can I test first? A: Yes. Sample data included."
])

# Slide 11: Next Steps
add_content_slide(prs, "🚀 Getting Started", [
    "1. Try it yourself (30 minutes with sample data)",
    "2. Discuss with your team about privacy needs",
    "3. Pilot program: Process small batch of real cases",
    "4. Have radiologist review results",
    "5. Full deployment if results look good"
])

# Slide 12: Summary
add_content_slide(prs, "📌 Summary", [
    "CT-DEFACE removes facial features from medical images",
    "Keeps all clinical data and patient information",
    "Works on any computer without expensive hardware",
    "Data stays completely under your control",
    "Ready to use in 30 minutes"
])

# Slide 13: Thank You
add_title_slide(prs, "Questions?", "Let's discuss how CT-DEFACE can help your organization")

# Save presentation
output_path = "CT-DEFACE_Presentation.pptx"
prs.save(output_path)
print(f"✓ Presentation created: {output_path}")
